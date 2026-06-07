from pptx import Presentation

def extract_ppt_slides(path):

    prs = Presentation(path)

    slides_data = []

    for slide in prs.slides:

        slide_text = ""

        for shape in slide.shapes:

            if hasattr(shape, "text"):

                slide_text += shape.text + "\n"

        slides_data.append(
            slide_text.strip()
        )

    return slides_data